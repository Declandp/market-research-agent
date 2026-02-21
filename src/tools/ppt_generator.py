"""
PowerPoint Report Generator
============================
Converts the markdown competitive intelligence report into a
professional PowerPoint presentation with selectable color themes
and template styles.
"""

import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ─────────────────────────────────────────────
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MED_TEXT = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT_TEXT = RGBColor(0x8C, 0x8C, 0x8C)
RULE_COLOR = RGBColor(0xD4, 0xD4, 0xD4)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Color Themes ──────────────────────────────────────────
THEMES = {
    "classic_green": {
        "title_bg":     RGBColor(0x00, 0x3B, 0x2D),
        "title_accent": RGBColor(0x00, 0x8A, 0x68),
        "title_sub":    RGBColor(0x7F, 0xBF, 0xAD),
        "title_muted":  RGBColor(0x5A, 0x9A, 0x8A),
        "primary":      RGBColor(0x00, 0x6B, 0x4F),
        "positive":     RGBColor(0x00, 0x7A, 0x87),  # teal
        "negative":     RGBColor(0xC4, 0x2B, 0x2B),
        "warning":      RGBColor(0xB8, 0x6E, 0x00),
        "alt":          RGBColor(0x4A, 0x5A, 0x6A),
        "card_bg":      RGBColor(0xF5, 0xF7, 0xF6),
        "table_header": RGBColor(0x00, 0x3B, 0x2D),
        "table_alt":    RGBColor(0xF5, 0xF7, 0xF6),
    },
    "navy_blue": {
        "title_bg":     RGBColor(0x1B, 0x2A, 0x4A),
        "title_accent": RGBColor(0x3B, 0x82, 0xF6),
        "title_sub":    RGBColor(0x93, 0xB5, 0xED),
        "title_muted":  RGBColor(0x64, 0x8A, 0xC7),
        "primary":      RGBColor(0x25, 0x63, 0xEB),
        "positive":     RGBColor(0x05, 0x96, 0x69),
        "negative":     RGBColor(0xDC, 0x26, 0x26),
        "warning":      RGBColor(0xD9, 0x77, 0x06),
        "alt":          RGBColor(0x4F, 0x46, 0xE5),
        "card_bg":      RGBColor(0xEF, 0xF6, 0xFF),
        "table_header": RGBColor(0x1B, 0x2A, 0x4A),
        "table_alt":    RGBColor(0xF0, 0xF4, 0xF8),
    },
    "charcoal": {
        "title_bg":     RGBColor(0x1C, 0x1C, 0x1C),
        "title_accent": RGBColor(0xC9, 0xA2, 0x27),
        "title_sub":    RGBColor(0xBD, 0xB7, 0x97),
        "title_muted":  RGBColor(0x8A, 0x86, 0x74),
        "primary":      RGBColor(0xC9, 0xA2, 0x27),
        "positive":     RGBColor(0x0D, 0x94, 0x88),
        "negative":     RGBColor(0xBE, 0x12, 0x3C),
        "warning":      RGBColor(0xB8, 0x6E, 0x00),
        "alt":          RGBColor(0x64, 0x74, 0x8B),
        "card_bg":      RGBColor(0xF5, 0xF5, 0xF0),
        "table_header": RGBColor(0x1C, 0x1C, 0x1C),
        "table_alt":    RGBColor(0xF5, 0xF5, 0xF0),
    },
}

# ── Template Styles ───────────────────────────────────────
STYLES = {
    "consulting": {
        "heading_font": "Georgia",
        "body_font":    "Calibri",
        "use_cards":    False,
        "use_dividers": True,
    },
    "corporate": {
        "heading_font": "Calibri",
        "body_font":    "Calibri",
        "use_cards":    True,
        "use_dividers": False,
    },
    "minimal": {
        "heading_font": "Georgia",
        "body_font":    "Calibri",
        "use_cards":    False,
        "use_dividers": False,
    },
}

AVAILABLE_THEMES = list(THEMES.keys())
AVAILABLE_STYLES = list(STYLES.keys())


# ── Helpers ───────────────────────────────────────────────

def _get_blank_layout(prs):
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def _clean(text: str) -> str:
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    return text.strip()


def _bg(slide, color=BG_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _rrect(slide, left, top, w, h, color, border_color=RULE_COLOR):
    """Rounded rectangle card."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.color.rgb = border_color
    s.line.width = Pt(0.75)
    return s


def _line(slide, left, top, width):
    return _rect(slide, left, top, width, Inches(0.015), RULE_COLOR)


def _text(slide, left, top, w, h, text, size=16, color=DARK_TEXT,
          bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = _clean(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def _bullets(tf, items, size=13, color=DARK_TEXT, label_color=None, body_font="Calibri"):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clean = item.strip().lstrip("*-•")
        clean = re.sub(r'^\d+[\.\)]\s*', '', clean)
        clean = _clean(clean)

        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = body_font
        p.space_before = Pt(2)
        p.space_after = Pt(6)

        if label_color and ":" in clean:
            parts = clean.split(":", 1)
            r1 = p.add_run()
            r1.text = parts[0] + ":"
            r1.font.size = Pt(size)
            r1.font.color.rgb = label_color
            r1.font.bold = True
            r1.font.name = body_font
            if len(parts) > 1:
                r2 = p.add_run()
                r2.text = " " + parts[1].strip()
                r2.font.size = Pt(size)
                r2.font.color.rgb = color
                r2.font.name = body_font
        else:
            p.text = clean


def _section_header(slide, text, T, S):
    """Section header — adapts to style."""
    hfont = S["heading_font"]
    primary = T["primary"]

    if S["use_cards"]:
        # Corporate: colored bar across top
        _rect(slide, Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.04), primary)
        _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
              text, size=24, color=BLACK, bold=True, font=hfont)
    else:
        # Consulting/Minimal: thin accent bar + title + horizontal rule
        _rect(slide, Inches(0.8), Inches(0.45), Inches(0.05), Inches(0.4), primary)
        _text(slide, Inches(1.0), Inches(0.35), Inches(10), Inches(0.6),
              text, size=24, color=BLACK, bold=True, font=hfont)
        if S["use_dividers"]:
            _line(slide, Inches(0.8), Inches(1.0), Inches(11.7))


def _col_separator(slide, S):
    """Add vertical column divider if style uses dividers."""
    if S["use_dividers"]:
        _rect(slide, Inches(6.65), Inches(1.3), Inches(0.015), Inches(5.5), RULE_COLOR)


def _card_or_nothing(slide, left, top, w, h, T, S):
    """Add card background if style uses cards."""
    if S["use_cards"]:
        _rrect(slide, left, top, w, h, T["card_bg"])


# ── Parsing ───────────────────────────────────────────────

def _parse_sections(report_text: str) -> dict:
    sections = {}
    current_key = "intro"
    current_content = []
    for line in report_text.split("\n"):
        if line.startswith("## "):
            if current_content:
                sections[current_key] = "\n".join(current_content)
            current_key = re.sub(r'^#+\s*', '', line).strip().rstrip(".")
            current_content = []
        elif line.startswith("# ") and current_key == "intro":
            sections["title"] = re.sub(r'^#+\s*', '', line).strip()
        else:
            current_content.append(line)
    if current_content:
        sections[current_key] = "\n".join(current_content)
    return sections


def _get_bullets(text: str) -> list[str]:
    bullets = []
    for line in text.split("\n"):
        s = line.strip()
        if s and (s.startswith(("-", "*", "•")) or re.match(r'^\d+[\.\)]', s)):
            c = re.sub(r'^[-*•]\s*', '', s)
            c = re.sub(r'^\d+[\.\)]\s*', '', c)
            c = _clean(c)
            if c and len(c) > 3:
                bullets.append(c)
    return bullets


def _get_table(text: str):
    headers, rows = [], []
    for line in text.split("\n"):
        s = line.strip()
        if "|" in s and not re.match(r'^\|[\s\-:|]+\|$', s):
            cells = [c.strip() for c in s.split("|") if c.strip()]
            if not headers:
                headers = cells
            else:
                rows.append(cells)
    return headers, rows


def _parse_swot(section_text: str) -> dict:
    swot = {"Strengths": [], "Weaknesses": [], "Opportunities": [], "Threats": []}
    current = None
    for line in section_text.split("\n"):
        stripped = line.strip()
        lower = stripped.lower()
        matched = None
        for key in swot:
            if key.lower() in lower:
                matched = key
                break
        if matched:
            cm = re.search(r'(?:strengths|weaknesses|opportunities|threats)\s*:?\s*:(.+)',
                           stripped, re.IGNORECASE)
            if cm:
                items = [i.strip().rstrip(".") for i in re.split(r'[,;]', cm.group(1)) if i.strip()]
                swot[matched].extend(items)
                current = matched
            elif ":" in stripped or stripped.startswith("#") or stripped.startswith("*"):
                current = matched
        elif current and stripped.startswith(("-", "*", "•")):
            clean = re.sub(r'^[-*•]\s*', '', stripped)
            clean = _clean(clean).rstrip(".")
            sub = None
            for key in swot:
                if clean.lower().startswith(key.lower()):
                    sub = key
                    break
            if sub:
                cp = clean.find(":")
                if cp > 0:
                    items = [i.strip().rstrip(".") for i in re.split(r'[,;]', clean[cp+1:]) if i.strip()]
                    swot[sub].extend(items)
            elif clean and len(clean) > 3:
                swot[current].append(clean)
    return swot


# ── Slide Builders ────────────────────────────────────────

def _slide_title(prs, company, competitors, date_str, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide, T["title_bg"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04), T["title_accent"])

    _text(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(0.6),
          "Competitive Intelligence Report", size=16,
          color=T["title_sub"], font=S["body_font"])

    _text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
          company, size=44, color=BG_WHITE, bold=True, font=S["heading_font"],
          align=PP_ALIGN.LEFT)

    _rect(slide, Inches(1.5), Inches(3.9), Inches(3), Inches(0.02), BG_WHITE)

    _text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
          f"Competitors: {', '.join(competitors)}", size=14,
          color=T["title_sub"], font=S["body_font"])

    _text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.4),
          date_str, size=12, color=T["title_muted"], font=S["body_font"])

    _text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.3),
          "CONFIDENTIAL", size=9, color=T["title_muted"],
          font=S["body_font"], bold=True)

    _rect(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04), T["title_accent"])


def _slide_exec_summary(prs, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, "Executive Summary", T, S)

    bullets_list = _get_bullets(section_text)
    if not bullets_list:
        bullets_list = [_clean(l.strip()) for l in section_text.split("\n")
                        if l.strip() and not l.strip().startswith("#")]

    bf = S["body_font"]
    lc = T["primary"]

    if len(bullets_list) > 5:
        mid = len(bullets_list) // 2

        _card_or_nothing(slide, Inches(0.8), Inches(1.2), Inches(5.7), Inches(5.8), T, S)
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(5.3), Inches(5.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, bullets_list[:mid], size=13, label_color=lc, body_font=bf)

        _col_separator(slide, S)

        _card_or_nothing(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(5.8), T, S)
        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.4), Inches(5.3), Inches(5.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, bullets_list[mid:mid+6], size=13, label_color=lc, body_font=bf)
    else:
        _card_or_nothing(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.8), T, S)
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(5.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, bullets_list[:8], size=14, label_color=lc, body_font=bf)


def _slide_market(prs, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, "Market Landscape", T, S)

    paragraphs = [p.strip() for p in section_text.split("\n")
                  if p.strip() and not p.strip().startswith("#")]
    blist = _get_bullets(section_text)
    bf = S["body_font"]
    hf = S["heading_font"]

    body = []
    for p in paragraphs:
        if not p.startswith(("-", "*", "•")) and not re.match(r'^\d+[\.\)]', p):
            c = _clean(p)
            if c and len(c) > 10:
                body.append(c)

    if body:
        _card_or_nothing(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(5.8), T, S)
        _text(slide, Inches(1.0), Inches(1.4), Inches(5.2), Inches(5.4),
              "\n\n".join(body[:4]), size=12, color=MED_TEXT, font=bf)

    if blist:
        _col_separator(slide, S)

        _card_or_nothing(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(5.8), T, S)
        _text(slide, Inches(7.0), Inches(1.3), Inches(5), Inches(0.4),
              "Key Trends", size=16, color=T["primary"], bold=True, font=hf)
        if S["use_dividers"]:
            _line(slide, Inches(7.0), Inches(1.75), Inches(5.5))

        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.3), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, blist[:6], size=12, label_color=T["primary"], body_font=bf)


def _slide_competitor(prs, name, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, f"Competitor: {name}", T, S)

    threat_match = re.search(r'[Tt]hreat\s*[Ll]evel.*?(\d+)/10', section_text)
    threat_level = int(threat_match.group(1)) if threat_match else 5
    badge_color = T["positive"] if threat_level <= 3 else (T["warning"] if threat_level <= 6 else T["negative"])
    _rect(slide, Inches(10.8), Inches(0.35), Inches(1.8), Inches(0.55), badge_color)
    _text(slide, Inches(10.8), Inches(0.38), Inches(1.8), Inches(0.5),
          f"Threat: {threat_level}/10", size=13, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER)

    blist = _get_bullets(section_text)
    if not blist:
        blist = [_clean(l.strip()) for l in section_text.split("\n")
                 if l.strip() and not l.strip().startswith("#")]

    mid = max(1, len(blist) // 2)
    bf = S["body_font"]
    hf = S["heading_font"]

    # Left
    _card_or_nothing(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(5.8), T, S)
    _text(slide, Inches(1.0), Inches(1.3), Inches(5), Inches(0.4),
          "Overview", size=16, color=T["primary"], bold=True, font=hf)
    if S["use_dividers"]:
        _line(slide, Inches(0.8), Inches(1.75), Inches(5.5))
    if blist[:mid]:
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.2), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, blist[:mid][:7], size=12, label_color=T["primary"], body_font=bf)

    _col_separator(slide, S)

    # Right
    _card_or_nothing(slide, Inches(6.8), Inches(1.2), Inches(5.6), Inches(5.8), T, S)
    _text(slide, Inches(7.0), Inches(1.3), Inches(5), Inches(0.4),
          "Strengths & Weaknesses", size=16, color=T["primary"], bold=True, font=hf)
    if S["use_dividers"]:
        _line(slide, Inches(7.0), Inches(1.75), Inches(5.5))
    if blist[mid:]:
        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.2), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, blist[mid:][:7], size=12, label_color=T["primary"], body_font=bf)


def _slide_matrix(prs, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, "Competitive Comparison", T, S)

    headers, rows = _get_table(section_text)
    bf = S["body_font"]

    if not headers or not rows:
        blist = _get_bullets(section_text)
        if blist:
            txBox = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.7))
            tf = txBox.text_frame
            tf.word_wrap = True
            _bullets(tf, blist[:10], size=13, label_color=T["primary"], body_font=bf)
        return

    nc = len(headers)
    nr = min(len(rows) + 1, 8)
    rows = rows[:nr - 1]

    tw = Inches(min(11.5, nc * 2.2))
    th = Inches(0.5 * nr)
    ts = slide.shapes.add_table(nr, nc, Inches(0.8), Inches(1.3), tw, th)
    table = ts.table

    cw = int(tw / nc)
    for col in table.columns:
        col.width = cw

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = _clean(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = T["table_header"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = BG_WHITE
            p.font.bold = True
            p.font.name = bf

    for r, rd in enumerate(rows):
        bg = BG_WHITE if r % 2 == 0 else T["table_alt"]
        for c in range(nc):
            cell = table.cell(r + 1, c)
            cell.text = _clean(rd[c]) if c < len(rd) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_TEXT
                p.font.name = bf


def _slide_swot(prs, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, "SWOT Analysis", T, S)

    swot = _parse_swot(section_text)
    bf = S["body_font"]
    hf = S["heading_font"]

    grid = [
        (Inches(0.8),  Inches(1.3), "Strengths",     T["primary"]),
        (Inches(6.8),  Inches(1.3), "Weaknesses",    T["negative"]),
        (Inches(0.8),  Inches(4.2), "Opportunities", T["positive"]),
        (Inches(6.8),  Inches(4.2), "Threats",        T["warning"]),
    ]

    for left, top, title, color in grid:
        items = swot.get(title, [])

        if S["use_cards"]:
            _rrect(slide, left, top, Inches(5.7), Inches(2.7), T["card_bg"])
        else:
            _rect(slide, left, top, Inches(5.7), Inches(2.7), T["card_bg"])

        _rect(slide, left, top, Inches(5.7), Inches(0.04), color)

        _text(slide, left + Inches(0.2), top + Inches(0.12),
              Inches(5), Inches(0.35), title,
              size=15, color=color, bold=True, font=hf)

        if items:
            txBox = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.55), Inches(5.3), Inches(2.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            _bullets(tf, items[:5], size=11, color=MED_TEXT, label_color=color, body_font=bf)
        else:
            _text(slide, left + Inches(0.2), top + Inches(0.7),
                  Inches(5), Inches(0.3), "No data available",
                  size=10, color=LIGHT_TEXT)


def _slide_opps_threats(prs, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, "Opportunities & Threats", T, S)

    opps, threats = [], []
    cur = None
    for line in section_text.split("\n"):
        s = line.strip().lower()
        if "opportunit" in s and (":" in s or s.startswith("#")):
            cur = "o"
        elif "threat" in s and (":" in s or s.startswith("#")):
            cur = "t"
        elif "gap" in s and (":" in s or s.startswith("#")):
            cur = None
        raw = line.strip()
        is_b = raw.startswith(("-", "*", "•")) or re.match(r'^\d+[\.\)]', raw)
        if cur and is_b:
            c = re.sub(r'^[-*•]\s*', '', raw)
            c = re.sub(r'^\d+[\.\)]\s*', '', c)
            c = _clean(c)
            if c and len(c) > 3:
                (opps if cur == "o" else threats).append(c)

    bf = S["body_font"]
    hf = S["heading_font"]

    # Left: Opportunities
    _card_or_nothing(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(5.8), T, S)
    _text(slide, Inches(1.0), Inches(1.3), Inches(5.5), Inches(0.4),
          "Opportunities", size=16, color=T["positive"], bold=True, font=hf)
    if S["use_dividers"]:
        _line(slide, Inches(0.8), Inches(1.75), Inches(5.5))
    if opps:
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.2), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, opps[:6], size=13, label_color=T["positive"], body_font=bf)

    _col_separator(slide, S)

    # Right: Threats
    _card_or_nothing(slide, Inches(6.8), Inches(1.2), Inches(5.6), Inches(5.8), T, S)
    _text(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
          "Threats", size=16, color=T["negative"], bold=True, font=hf)
    if S["use_dividers"]:
        _line(slide, Inches(7.0), Inches(1.75), Inches(5.5))
    if threats:
        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.2), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        _bullets(tf, threats[:6], size=13, label_color=T["negative"], body_font=bf)


def _slide_recommendations(prs, section_text, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide)
    _section_header(slide, "Strategic Recommendations", T, S)

    timeframes = [
        ("Immediate (30 Days)", T["primary"], []),
        ("Short-Term (3 Months)", T["positive"], []),
        ("Long-Term (12 Months)", T["alt"], []),
    ]

    idx = None
    for line in section_text.split("\n"):
        s = line.strip().lower()
        if "immediate" in s or "30 day" in s or "next 30" in s:
            idx = 0
        elif "short" in s or "3 month" in s or "next 3" in s:
            idx = 1
        elif "long" in s or "12 month" in s or "next 12" in s:
            idx = 2
        raw = line.strip()
        is_b = raw.startswith(("-", "*", "•")) or re.match(r'^\d+[\.\)]', raw)
        if idx is not None and is_b:
            c = re.sub(r'^[-*•]\s*', '', raw)
            c = re.sub(r'^\d+[\.\)]\s*', '', c)
            c = _clean(c)
            c = re.sub(r'\[.\]', '', c).strip()
            if c and len(c) > 3:
                timeframes[idx][2].append(c)

    bf = S["body_font"]
    hf = S["heading_font"]
    col_w = Inches(3.8)
    gap = Inches(0.15)

    for i, (label, color, items) in enumerate(timeframes):
        left = Inches(0.8) + (col_w + gap) * i

        _card_or_nothing(slide, left, Inches(1.2), col_w, Inches(5.8), T, S)
        _rect(slide, left, Inches(1.3), col_w, Inches(0.04), color)
        _text(slide, left + Inches(0.15), Inches(1.45), col_w - Inches(0.3),
              Inches(0.4), label, size=14, color=color, bold=True, font=hf)
        if S["use_dividers"]:
            _line(slide, left, Inches(1.9), col_w)

        if items:
            txBox = slide.shapes.add_textbox(
                left + Inches(0.15), Inches(2.1), col_w - Inches(0.3), Inches(4.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            _bullets(tf, items[:5], size=12, color=MED_TEXT, label_color=color, body_font=bf)

    # Column dividers
    if S["use_dividers"]:
        _rect(slide, Inches(0.8) + col_w + gap / 2, Inches(1.3),
              Inches(0.015), Inches(5.5), RULE_COLOR)
        _rect(slide, Inches(0.8) + (col_w + gap) * 2 - gap / 2, Inches(1.3),
              Inches(0.015), Inches(5.5), RULE_COLOR)


def _slide_closing(prs, T, S):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _bg(slide, T["title_bg"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04), T["title_accent"])

    _text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8),
          "Thank You", size=44, color=BG_WHITE, bold=True, font=S["heading_font"],
          align=PP_ALIGN.LEFT)

    _rect(slide, Inches(1.5), Inches(3.5), Inches(2.5), Inches(0.02), BG_WHITE)

    _text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5),
          "AI Market Research Agent", size=18,
          color=T["title_sub"], bold=True, font=S["body_font"])

    _text(slide, Inches(1.5), Inches(4.4), Inches(10), Inches(0.4),
          "Multi-Agent Intelligence System", size=14,
          color=T["title_muted"], font=S["body_font"])

    _text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.3),
          "CONFIDENTIAL", size=9, color=T["title_muted"],
          font=S["body_font"], bold=True)

    _rect(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04), T["title_accent"])


# ── Main Generator ────────────────────────────────────────

def generate_pptx(report_text: str, company: str, competitors: list[str],
                  output_path: str, theme: str = "classic_green",
                  style: str = "consulting") -> str:
    """
    Generate a PowerPoint presentation from a markdown report.

    Args:
        report_text: Markdown report content
        company: Company name
        competitors: List of competitor names
        output_path: Path to save the .pptx file
        theme: Color theme - "classic_green", "navy_blue", or "charcoal"
        style: Template style - "consulting", "corporate", or "minimal"
    """
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    T = THEMES.get(theme, THEMES["classic_green"])
    S = STYLES.get(style, STYLES["consulting"])

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    date_str = datetime.now().strftime("%B %d, %Y")
    sections = _parse_sections(report_text)

    _slide_title(prs, company, competitors, date_str, T, S)

    for key in sections:
        if "executive" in key.lower() or "summary" in key.lower():
            _slide_exec_summary(prs, sections[key], T, S)
            break

    for key in sections:
        if "market" in key.lower() and "landscape" in key.lower():
            _slide_market(prs, sections[key], T, S)
            break

    for key in sections:
        if "competitor" in key.lower() and "profile" in key.lower():
            st = sections[key]
            for comp in competitors:
                if comp.lower() in st.lower():
                    _slide_competitor(prs, comp, st, T, S)
                    break
            else:
                _slide_competitor(prs, ", ".join(competitors), st, T, S)
            break

    for key in sections:
        if "competitive" in key.lower() and "analy" in key.lower():
            _slide_matrix(prs, sections[key], T, S)
            if "swot" in sections[key].lower():
                _slide_swot(prs, sections[key], T, S)
            break

    for key in sections:
        if "swot" in key.lower():
            _slide_swot(prs, sections[key], T, S)
            break

    for key in sections:
        if "opportunit" in key.lower() or "threat" in key.lower():
            _slide_opps_threats(prs, sections[key], T, S)
            break

    for key in sections:
        if "recommend" in key.lower() or "strateg" in key.lower():
            _slide_recommendations(prs, sections[key], T, S)
            break

    _slide_closing(prs, T, S)

    prs.save(output_path)
    return output_path
